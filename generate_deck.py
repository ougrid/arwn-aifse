"""Generate the presentation deck for CS Demand Forecasting & Shift Scheduling.

Brand CI: Aura Bangkok Clinic
- Primary: #00C3F7 (bright cyan/sky blue)
- Secondary: #0097D6 (deeper blue for accents)
- Dark: #1A1A2E (near-black for text)
- White: #FFFFFF (text on colored backgrounds)
- Accent: #FF6B6B (warm accent for highlights/warnings)
- Success: #2ECC71 (green for good metrics)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──────────────────────────────────────────────────────
CYAN = RGBColor(0x00, 0xC3, 0xF7)  # Primary brand blue
DEEP_BLUE = RGBColor(0x00, 0x97, 0xD6)  # Darker accent
DARK = RGBColor(0x1A, 0x1A, 0x2E)  # Dark text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xFA)
MEDIUM_GRAY = RGBColor(0x8E, 0x8E, 0xA3)
ACCENT_RED = RGBColor(0xFF, 0x6B, 0x6B)
SUCCESS = RGBColor(0x2E, 0xCC, 0x71)
AMBER = RGBColor(0xF3, 0x9C, 0x12)

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, corner_radius=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # no border
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=DARK,
    bold=False,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(
    slide, left, top, width, height, items, font_size=16, color=DARK, line_spacing=1.5
):
    """Add a bullet point list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * (line_spacing - 1))
        p.level = 0
    return txBox


def add_metric_card(
    slide,
    left,
    top,
    width,
    height,
    label,
    value,
    bg_color=WHITE,
    value_color=CYAN,
    label_color=MEDIUM_GRAY,
):
    """Add a metric card with label and value."""
    card = add_shape(slide, left, top, width, height, bg_color, corner_radius=0.05)

    # Value
    add_textbox(
        slide,
        left + Inches(0.2),
        top + Inches(0.15),
        width - Inches(0.4),
        Inches(0.6),
        value,
        font_size=28,
        color=value_color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    # Label
    add_textbox(
        slide,
        left + Inches(0.2),
        top + Inches(0.7),
        width - Inches(0.4),
        Inches(0.4),
        label,
        font_size=11,
        color=label_color,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )
    return card


def add_section_badge(slide, left, top, text, bg_color=CYAN):
    """Add a small colored badge/tag."""
    shape = add_shape(
        slide, left, top, Inches(1.8), Inches(0.35), bg_color, corner_radius=0.15
    )
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(10)
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Calibri"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    return shape


def add_footer(
    slide, text="Aura Bangkok Clinic  |  CS Demand Forecasting & Shift Scheduling"
):
    """Add a subtle footer to a slide."""
    add_textbox(
        slide,
        Inches(0.5),
        Inches(7.0),
        Inches(12),
        Inches(0.3),
        text,
        font_size=9,
        color=MEDIUM_GRAY,
        alignment=PP_ALIGN.CENTER,
    )


def add_top_bar(slide, color=CYAN, height=Inches(0.06)):
    """Add thin colored bar at top of slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def make_slide_header(slide, title, subtitle=None):
    """Standard slide header with top bar."""
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.4),
        Inches(11),
        Inches(0.7),
        title,
        font_size=32,
        color=DARK,
        bold=True,
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.8),
            Inches(1.0),
            Inches(11),
            Inches(0.5),
            subtitle,
            font_size=16,
            color=MEDIUM_GRAY,
        )
    add_footer(slide)


# ── Slide Builders ────────────────────────────────────────────────────


def slide_title(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, CYAN)

    # Main title
    add_textbox(
        slide,
        Inches(1),
        Inches(1.8),
        Inches(11),
        Inches(1.2),
        "CS Demand Forecasting\n& Shift Scheduling",
        font_size=44,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.LEFT,
    )

    # Subtitle
    add_textbox(
        slide,
        Inches(1),
        Inches(3.4),
        Inches(11),
        Inches(0.5),
        "AI-Powered Staffing Optimization for April 2026",
        font_size=20,
        color=WHITE,
        bold=False,
        alignment=PP_ALIGN.LEFT,
    )

    # Divider line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.1), Inches(3), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()

    # Author info
    add_textbox(
        slide,
        Inches(1),
        Inches(4.4),
        Inches(11),
        Inches(0.4),
        "Technical Assessment  |  AI & Full-Stack Engineer",
        font_size=16,
        color=WHITE,
        alignment=PP_ALIGN.LEFT,
    )

    # Brand badge
    add_textbox(
        slide,
        Inches(1),
        Inches(5.8),
        Inches(11),
        Inches(0.4),
        "Aura Bangkok Clinic",
        font_size=14,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.LEFT,
    )


def slide_problem(prs):
    """Slide 2: Problem Statement."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, "The Challenge", "What we're solving")

    # Problem cards
    cards = [
        ("50 Agents", "10 Senior + 40 Junior\n8 English-capable", CYAN),
        ("4 Shifts/Day", "Morning · Afternoon\nEvening · Night", DEEP_BLUE),
        ("30 Days", "April 2026\n6 leave days each", CYAN),
        ("2 KPI Targets", "CSAT ≥ 4.0\nWait ≤ 60 seconds", DEEP_BLUE),
    ]

    for i, (title, desc, color) in enumerate(cards):
        left = Inches(0.8 + i * 3.1)
        card = add_shape(
            slide, left, Inches(1.8), Inches(2.8), Inches(2.2), color, 0.05
        )
        add_textbox(
            slide,
            left + Inches(0.2),
            Inches(2.0),
            Inches(2.4),
            Inches(0.5),
            title,
            font_size=22,
            color=WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left + Inches(0.2),
            Inches(2.6),
            Inches(2.4),
            Inches(1.0),
            desc,
            font_size=14,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
        )

    # Constraints summary
    add_textbox(
        slide,
        Inches(0.8),
        Inches(4.4),
        Inches(11),
        Inches(0.4),
        "Business Constraints",
        font_size=20,
        color=DARK,
        bold=True,
    )

    constraints = [
        "6 hard constraints: leave allowance, pre-selected leave, night rest, minimum staffing, daily leave cap, valid assignments",
        "2 soft objectives: night shift fairness across agents, minimize overstaffing",
        "Goal: Find an optimal schedule respecting ALL constraints while maximizing service quality",
    ]
    add_bullet_list(
        slide,
        Inches(0.8),
        Inches(4.9),
        Inches(11.5),
        Inches(2.0),
        constraints,
        font_size=14,
        color=DARK,
    )


def slide_approach(prs):
    """Slide 3: Solution approach overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(
        slide,
        "Solution Architecture",
        "Two-stage pipeline: Predict → Optimize → Validate",
    )

    # Stage boxes
    stages = [
        (
            "1",
            "FORECAST",
            "Predict ticket volume\nper shift per day",
            "GradientBoosting\nRegressor",
            CYAN,
        ),
        (
            "2",
            "STAFFING",
            "Determine minimum\nSr/Jr/Eng per shift",
            "Quality Models\n+ Grid Search",
            DEEP_BLUE,
        ),
        (
            "3",
            "SCHEDULE",
            "Assign 50 agents\nto shifts × 30 days",
            "Google OR-Tools\nCP-SAT Solver",
            CYAN,
        ),
        (
            "4",
            "EVALUATE",
            "Validate constraints\n& project quality",
            "Fairness + CSAT\n+ Remediation",
            DEEP_BLUE,
        ),
    ]

    for i, (num, title, desc, tech, color) in enumerate(stages):
        left = Inches(0.6 + i * 3.15)
        # Card background
        add_shape(slide, left, Inches(1.8), Inches(2.8), Inches(3.2), color, 0.05)
        # Step number
        add_textbox(
            slide,
            left + Inches(0.2),
            Inches(1.9),
            Inches(2.4),
            Inches(0.5),
            f"Stage {num}",
            font_size=12,
            color=WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        # Title
        add_textbox(
            slide,
            left + Inches(0.2),
            Inches(2.3),
            Inches(2.4),
            Inches(0.5),
            title,
            font_size=20,
            color=WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        # Description
        add_textbox(
            slide,
            left + Inches(0.2),
            Inches(2.9),
            Inches(2.4),
            Inches(0.8),
            desc,
            font_size=13,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
        )
        # Tech badge
        tech_box = add_shape(
            slide,
            left + Inches(0.3),
            Inches(3.9),
            Inches(2.2),
            Inches(0.7),
            WHITE,
            0.08,
        )
        add_textbox(
            slide,
            left + Inches(0.3),
            Inches(3.95),
            Inches(2.2),
            Inches(0.6),
            tech,
            font_size=11,
            color=color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Arrow between stages
        if i < 3:
            arrow_left = left + Inches(2.85)
            add_textbox(
                slide,
                arrow_left,
                Inches(3.0),
                Inches(0.3),
                Inches(0.5),
                "→",
                font_size=28,
                color=MEDIUM_GRAY,
                bold=True,
                alignment=PP_ALIGN.CENTER,
            )

    # Tech stack bar
    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.5),
        Inches(11),
        Inches(0.4),
        "Tech Stack:  Python 3.11  ·  scikit-learn  ·  Google OR-Tools  ·  Streamlit  ·  Plotly  ·  pandas  ·  pytest",
        font_size=13,
        color=MEDIUM_GRAY,
        alignment=PP_ALIGN.CENTER,
    )


def slide_forecasting(prs):
    """Slide 4: Demand Forecasting deep dive."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(
        slide,
        "Task A: Demand Forecasting",
        "Two-stage approach — Volume Prediction → Staffing Optimization",
    )

    # Left column: Volume Model
    add_section_badge(slide, Inches(0.8), Inches(1.6), "STAGE 1: VOLUME")
    add_textbox(
        slide,
        Inches(0.8),
        Inches(2.1),
        Inches(5.5),
        Inches(0.4),
        "GradientBoostingRegressor",
        font_size=18,
        color=DARK,
        bold=True,
    )

    vol_details = [
        "Trained on 6 months of historical shift data (728 records)",
        "Features: day_of_week, is_weekend, cyclical month, shift ID, day_of_month",
        "Time-based validation split (last month held out)",
        "Predicts ticket volume per shift per day for April 2026",
    ]
    add_bullet_list(
        slide,
        Inches(0.8),
        Inches(2.5),
        Inches(5.5),
        Inches(2.0),
        vol_details,
        font_size=13,
    )

    # Volume metrics
    add_metric_card(
        slide,
        Inches(0.8),
        Inches(4.4),
        Inches(2.5),
        Inches(1.1),
        "Validation R²",
        "0.927",
        value_color=SUCCESS,
    )
    add_metric_card(
        slide,
        Inches(3.6),
        Inches(4.4),
        Inches(2.5),
        Inches(1.1),
        "Validation MAE",
        "19.7 tickets",
        value_color=CYAN,
    )

    # Right column: Staffing Model
    add_section_badge(slide, Inches(7), Inches(1.6), "STAGE 2: STAFFING")
    add_textbox(
        slide,
        Inches(7),
        Inches(2.1),
        Inches(5.5),
        Inches(0.4),
        "Quality Models + Grid Search",
        font_size=18,
        color=DARK,
        bold=True,
    )

    staff_details = [
        "CSAT Regression: volume + senior + junior + english → CSAT score",
        "Wait Time Regression: same features → avg wait seconds",
        "Vectorized grid search: find minimum staff meeting both targets",
        "Historical floors (p10) prevent model extrapolation errors",
        "Feasibility caps for scarce pools (8 English agents, ~6.4 avail/day)",
    ]
    add_bullet_list(
        slide,
        Inches(7),
        Inches(2.5),
        Inches(5.5),
        Inches(2.4),
        staff_details,
        font_size=13,
    )

    # Quality metrics
    add_metric_card(
        slide,
        Inches(7),
        Inches(4.4),
        Inches(2.5),
        Inches(1.1),
        "CSAT Model R²",
        "0.818",
        value_color=SUCCESS,
    )
    add_metric_card(
        slide,
        Inches(9.8),
        Inches(4.4),
        Inches(2.5),
        Inches(1.1),
        "Wait Model R²",
        "0.963",
        value_color=SUCCESS,
    )

    # Why GBR callout
    add_shape(
        slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.0), LIGHT_GRAY, 0.03
    )
    add_textbox(
        slide,
        Inches(1.0),
        Inches(5.85),
        Inches(11.3),
        Inches(0.9),
        "Why GradientBoosting over LSTM/Prophet?\n"
        "→ Tabular data with few features (not sequential), small dataset (728 rows), "
        "interpretable feature importance, fast training, strong out-of-box performance on structured data. "
        "LSTM would need 10-100x more data and offer no interpretability benefit.",
        font_size=12,
        color=DARK,
    )


def slide_scheduling(prs):
    """Slide 5: Scheduling Engine."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(
        slide,
        "Task B: Shift Scheduling Engine",
        "Constraint Programming with Google OR-Tools CP-SAT",
    )

    # Left: Hard Constraints
    add_shape(
        slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(4.4), LIGHT_GRAY, 0.03
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.8),
        Inches(5.4),
        Inches(0.4),
        "Hard Constraints (must satisfy)",
        font_size=18,
        color=ACCENT_RED,
        bold=True,
    )

    hard = [
        "1. Exactly one assignment per agent per day (shift or leave)",
        "2. Each agent gets exactly 6 leave days per month",
        "3. All 3 pre-selected leave dates honored",
        "4. Night shift → mandatory rest day (counts as leave)",
        "5. Minimum staffing per shift: Senior, Junior, English",
        "6. Daily leave cap: max 19 agents on leave",
    ]
    add_bullet_list(
        slide,
        Inches(0.8),
        Inches(2.3),
        Inches(5.4),
        Inches(3.5),
        hard,
        font_size=13,
        color=DARK,
    )

    # Right: Soft Constraints
    add_shape(
        slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.4), LIGHT_GRAY, 0.03
    )
    add_textbox(
        slide,
        Inches(7.0),
        Inches(1.8),
        Inches(5.4),
        Inches(0.4),
        "Soft Constraints (optimize)",
        font_size=18,
        color=AMBER,
        bold=True,
    )

    soft = [
        "1. Night shift fairness — minimize variance across agents (w=5)",
        "2. Overstaffing penalty — minimize excess beyond minimums (w=1)",
        "3. Agent preferences — honor shift likes/dislikes (w=1)",
    ]
    add_bullet_list(
        slide,
        Inches(7.0),
        Inches(2.3),
        Inches(5.4),
        Inches(2.0),
        soft,
        font_size=13,
        color=DARK,
    )

    # Model size
    add_textbox(
        slide,
        Inches(7.0),
        Inches(3.8),
        Inches(5.4),
        Inches(0.4),
        "Model Complexity",
        font_size=16,
        color=DARK,
        bold=True,
    )
    model_stats = [
        "Decision variables: 50 agents × 30 days × 5 options = 7,500",
        "Solver: CP-SAT (Constraint Programming – SAT)",
        "Solution time: ~60s → FEASIBLE (119/120 targets met)",
    ]
    add_bullet_list(
        slide,
        Inches(7.0),
        Inches(4.2),
        Inches(5.4),
        Inches(1.5),
        model_stats,
        font_size=13,
        color=DARK,
    )

    # Why CP-SAT callout
    add_shape(slide, Inches(0.6), Inches(6.3), Inches(12.0), Inches(0.8), CYAN, 0.03)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(6.35),
        Inches(11.6),
        Inches(0.7),
        "Why CP-SAT over heuristics/genetic algorithms?\n"
        "→ Guaranteed constraint satisfaction, provable optimality bounds, "
        "declarative modeling (add constraints without rewriting logic), "
        "scales well to 7,500 variables, industry-standard for shift scheduling.",
        font_size=12,
        color=WHITE,
    )


def slide_results(prs):
    """Slide 6: Results summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, "Results", "Full pipeline performance")

    # Top metric row
    metrics = [
        ("Avg CSAT", "4.27", "Target ≥ 4.0", SUCCESS),
        ("Avg Wait", "21.6s", "Target ≤ 60s", SUCCESS),
        ("Targets Met", "119/120", "99.2% of shifts", SUCCESS),
        ("Hard Errors", "0", "All satisfied", SUCCESS),
        ("Pipeline Time", "~71s", "End-to-end", CYAN),
    ]

    for i, (label, value, detail, color) in enumerate(metrics):
        left = Inches(0.6 + i * 2.5)
        add_shape(slide, left, Inches(1.6), Inches(2.2), Inches(1.6), LIGHT_GRAY, 0.05)
        add_textbox(
            slide,
            left + Inches(0.1),
            Inches(1.7),
            Inches(2.0),
            Inches(0.6),
            value,
            font_size=32,
            color=color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left + Inches(0.1),
            Inches(2.3),
            Inches(2.0),
            Inches(0.3),
            label,
            font_size=13,
            color=DARK,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left + Inches(0.1),
            Inches(2.65),
            Inches(2.0),
            Inches(0.3),
            detail,
            font_size=10,
            color=MEDIUM_GRAY,
            alignment=PP_ALIGN.CENTER,
        )

    # Staffing table header
    add_textbox(
        slide,
        Inches(0.8),
        Inches(3.6),
        Inches(11),
        Inches(0.4),
        "Average Daily Staffing by Shift",
        font_size=18,
        color=DARK,
        bold=True,
    )

    # Table-like layout
    headers = ["Shift", "Senior", "Junior", "English", "Total", "CSAT", "Wait (s)"]
    rows = [
        ["Morning (06-14)", "5.0", "10.3", "4.6", "15.3", "4.29", "18.1"],
        ["Afternoon (14-20)", "3.5", "8.2", "3.8", "11.7", "4.25", "23.4"],
        ["Evening (20-00)", "1.7", "3.8", "1.3", "5.5", "4.30", "19.8"],
        ["Night (00-06)", "0.5", "2.0", "0.7", "2.5", "4.22", "25.2"],
    ]

    # Header row
    col_widths = [2.2, 1.2, 1.2, 1.2, 1.2, 1.2, 1.2]
    y = Inches(4.1)
    x_start = Inches(1.5)
    # Header background first
    add_shape(slide, x_start, y, Inches(sum(col_widths)), Inches(0.35), CYAN)
    # Header text on top
    for j, (header, w) in enumerate(zip(headers, col_widths)):
        add_textbox(
            slide,
            x_start + Inches(sum(col_widths[:j])),
            y,
            Inches(w),
            Inches(0.35),
            header,
            font_size=12,
            color=WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

    # Data rows
    for i, row in enumerate(rows):
        y_row = Inches(4.5 + i * 0.4)
        bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
        add_shape(
            slide, x_start, y_row, Inches(sum(col_widths)), Inches(0.38), bg_color
        )
        for j, (val, w) in enumerate(zip(row, col_widths)):
            add_textbox(
                slide,
                x_start + Inches(sum(col_widths[:j])),
                y_row,
                Inches(w),
                Inches(0.38),
                val,
                font_size=12,
                color=DARK,
                alignment=PP_ALIGN.CENTER,
            )


def slide_decisions(prs):
    """Slide 7: Design decisions & trade-offs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(
        slide, "Design Decisions & Trade-offs", "Why we chose what we chose"
    )

    decisions = [
        (
            "GradientBoosting over LSTM/Prophet",
            "Tabular data (not sequential), only 728 rows, "
            "strong OOB performance, interpretable features, <1s training. "
            "LSTM needs 10-100x more data for marginal improvement.",
            "FORECASTING",
        ),
        (
            "CP-SAT over Heuristics/GA",
            "Guaranteed constraint satisfaction, provable bounds, "
            "declarative modeling (easy to add/modify constraints), "
            "scales well for 7,500 decision variables. Greedy/GA cannot guarantee feasibility.",
            "SCHEDULING",
        ),
        (
            "Two-Stage Forecasting",
            "Volume → Staffing decoupling: isolates forecasting error from optimization. "
            "Quality models bridge the gap (CSAT/wait as function of staffing). "
            "Historical floors prevent model extrapolation into unrealistic territory.",
            "ARCHITECTURE",
        ),
        (
            "Config-Driven Design",
            "All business rules in SchedulingConfig dataclass — weights, thresholds, "
            "solver params. Zero magic numbers in logic. "
            "Managers can tune via dashboard sliders without code changes.",
            "PRODUCTION",
        ),
    ]

    for i, (title, reasoning, badge) in enumerate(decisions):
        row = i // 2
        col = i % 2
        left = Inches(0.6 + col * 6.3)
        top = Inches(1.6 + row * 2.7)

        add_shape(slide, left, top, Inches(5.9), Inches(2.3), LIGHT_GRAY, 0.03)
        add_section_badge(
            slide, left + Inches(0.15), top + Inches(0.15), badge, DEEP_BLUE
        )
        add_textbox(
            slide,
            left + Inches(0.2),
            top + Inches(0.6),
            Inches(5.5),
            Inches(0.4),
            title,
            font_size=16,
            color=DARK,
            bold=True,
        )
        add_textbox(
            slide,
            left + Inches(0.2),
            top + Inches(1.0),
            Inches(5.5),
            Inches(1.2),
            reasoning,
            font_size=12,
            color=DARK,
        )


def slide_bonus_preferences(prs):
    """Slide 8: Bonus — Agent Preferences."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(
        slide,
        "Bonus: Agent Shift Preferences",
        "Soft constraint integration for employee satisfaction",
    )

    # Left: How it works
    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.6),
        Inches(5.5),
        Inches(0.4),
        "How It Works",
        font_size=20,
        color=DARK,
        bold=True,
    )

    pref_details = [
        "Each agent has preferences for each shift (1=love → 5=avoid)",
        "Generated based on role (seniors prefer day), language, + individual variation",
        "Integrated as soft objective: solver minimizes assigning agents to disliked shifts",
        "Preference cost = (preference score - 1) × weight in CP-SAT objective",
        "Balanced against fairness and staffing: satisfaction is a goal, not a guarantee",
    ]
    add_bullet_list(
        slide,
        Inches(0.8),
        Inches(2.1),
        Inches(5.5),
        Inches(2.5),
        pref_details,
        font_size=13,
    )

    # Right: Results
    add_textbox(
        slide,
        Inches(7),
        Inches(1.6),
        Inches(5.5),
        Inches(0.4),
        "Results",
        font_size=20,
        color=DARK,
        bold=True,
    )

    add_metric_card(
        slide,
        Inches(7),
        Inches(2.2),
        Inches(2.5),
        Inches(1.1),
        "Satisfaction",
        "82%",
        value_color=SUCCESS,
    )
    add_metric_card(
        slide,
        Inches(9.8),
        Inches(2.2),
        Inches(2.5),
        Inches(1.1),
        "Preferred Shifts",
        "1,083/1,200",
        value_color=CYAN,
    )

    add_textbox(
        slide,
        Inches(7),
        Inches(3.6),
        Inches(5.5),
        Inches(0.4),
        "Production Value",
        font_size=16,
        color=DARK,
        bold=True,
    )

    prod_details = [
        "Reduces burnout by respecting individual scheduling preferences",
        "Improves agent retention — a key metric for CS operations",
        "Tunable weight: managers can trade preference vs. fairness",
        "Dashboard tab shows per-agent satisfaction with heatmap visualization",
    ]
    add_bullet_list(
        slide,
        Inches(7),
        Inches(4.0),
        Inches(5.5),
        Inches(2.0),
        prod_details,
        font_size=13,
    )

    # Preference matrix placeholder
    add_shape(
        slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.0), LIGHT_GRAY, 0.03
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.5),
        Inches(5.5),
        Inches(0.5),
        "[Screenshot: Preference Heatmap from Dashboard]",
        font_size=14,
        color=MEDIUM_GRAY,
        alignment=PP_ALIGN.CENTER,
    )


def slide_bonus_fairness(prs):
    """Slide 9: Bonus — Fairness & Remediation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(
        slide,
        "Bonus: Fairness Metrics & Remediation",
        "Quantitative equity measurement + actionable suggestions",
    )

    # Left: Fairness
    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.6),
        Inches(5.5),
        Inches(0.4),
        "Fairness Scoring",
        font_size=20,
        color=DARK,
        bold=True,
    )

    fairness_details = [
        "Gini coefficient: 0 = everyone equal, 1 = one person gets all",
        "Night shift Gini: measures equity in unpopular shift distribution",
        "Workload Gini: measures equity in total working days",
        "Shift entropy (Shannon): diversity of shift assignments per agent",
        "Composite score (0-100): 40% night + 30% workload + 30% variety",
        "Letter grade: A (≥90) → F (<60); sparsity-adjusted for realism",
    ]
    add_bullet_list(
        slide,
        Inches(0.8),
        Inches(2.1),
        Inches(5.5),
        Inches(3.0),
        fairness_details,
        font_size=13,
    )

    # Right: Remediation
    add_textbox(
        slide,
        Inches(7),
        Inches(1.6),
        Inches(5.5),
        Inches(0.4),
        "Constraint Remediation",
        font_size=20,
        color=DARK,
        bold=True,
    )

    rem_details = [
        "Each violation type → specific, actionable suggestion",
        "Grouped by category with affected days listed",
        'e.g. "Move a senior from lower-demand shift" for staffing gaps',
        'e.g. "Swap with agent having complementary pattern" for continuity',
        "Dashboard shows expandable remediation per violation group",
        "Enables managers to make informed manual adjustments",
    ]
    add_bullet_list(
        slide,
        Inches(7),
        Inches(2.1),
        Inches(5.5),
        Inches(3.0),
        rem_details,
        font_size=13,
    )

    # Bottom: Why it matters
    add_shape(slide, Inches(0.6), Inches(5.4), Inches(12.0), Inches(1.2), CYAN, 0.03)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.5),
        Inches(11.6),
        Inches(1.0),
        "Why This Matters for Production\n"
        "→ Fairness metrics give managers visibility into equity issues before agent complaints. "
        "Remediation transforms a 'failed constraint' alert into a concrete action plan. "
        "Together, they close the feedback loop between automated scheduling and human oversight.",
        font_size=13,
        color=WHITE,
    )


def slide_production(prs):
    """Slide 10: Production readiness."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(
        slide,
        "Production Readiness",
        "Built for maintainability, not just to pass the test",
    )

    aspects = [
        (
            "Modular Architecture",
            "Clean separation: forecasting / scheduling / evaluation. "
            "Each module testable independently. New constraints = new function, not rewrite.",
            "📦",
        ),
        (
            "Config-Driven",
            "SchedulingConfig dataclass: all business rules, weights, thresholds in one place. "
            "Dashboard sliders map directly to config. Zero magic numbers in logic.",
            "⚙️",
        ),
        (
            "Comprehensive Testing",
            "55+ pytest tests: data integrity (8), forecasting (11), "
            "scheduling (13), bonus features (26). Module + integration coverage.",
            "✅",
        ),
        (
            "CI/CD Ready",
            "pyproject.toml with dependencies. Conventional commits (13 total). "
            "Tests run in <5s (unit) + ~71s (integration pipeline). Git-clean repo.",
            "🚀",
        ),
        (
            "Defensive Engineering",
            "Historical floors prevent model extrapolation. Feasibility caps for scarce agents. "
            "Windows encoding fix for cross-platform. Graceful degradation on solver timeout.",
            "🛡️",
        ),
        (
            "Interactive Dashboard",
            "7-tab Streamlit app: managers adjust parameters, see results, download CSV. "
            "Cached pipeline results for responsive interaction. No code changes needed.",
            "📊",
        ),
    ]

    for i, (title, desc, icon) in enumerate(aspects):
        row = i // 2
        col = i % 2
        left = Inches(0.6 + col * 6.3)
        top = Inches(1.6 + row * 1.7)

        add_shape(slide, left, top, Inches(5.9), Inches(1.4), LIGHT_GRAY, 0.03)
        add_textbox(
            slide,
            left + Inches(0.2),
            top + Inches(0.1),
            Inches(0.5),
            Inches(0.4),
            icon,
            font_size=20,
            color=DARK,
            alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left + Inches(0.6),
            top + Inches(0.1),
            Inches(5.0),
            Inches(0.35),
            title,
            font_size=15,
            color=DARK,
            bold=True,
        )
        add_textbox(
            slide,
            left + Inches(0.6),
            top + Inches(0.45),
            Inches(5.0),
            Inches(0.85),
            desc,
            font_size=11,
            color=DARK,
        )


def slide_dashboard(prs):
    """Slide 11: Dashboard preview with screenshot placeholder."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(
        slide,
        "Interactive Dashboard",
        "7 tabs: Forecast · Schedule · Summary · Constraints · Agents · Fairness · Preferences",
    )

    # Screenshot placeholder - large
    add_shape(
        slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.0), LIGHT_GRAY, 0.02
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(3.5),
        Inches(11.7),
        Inches(0.8),
        "[Insert Dashboard Screenshots Here]\nRun:  streamlit run app.py",
        font_size=20,
        color=MEDIUM_GRAY,
        alignment=PP_ALIGN.CENTER,
    )

    # Feature callouts at bottom
    features = [
        "Adjustable CSAT/wait targets",
        "Calendar schedule view",
        "Download CSV export",
        "Real-time constraint validation",
        "Preference heatmap",
        "Fairness grade & Gini",
        "Remediation suggestions",
    ]
    for i, feat in enumerate(features):
        left = Inches(0.5 + i * 1.8)
        add_shape(slide, left, Inches(6.6), Inches(1.6), Inches(0.4), CYAN, 0.1)
        add_textbox(
            slide,
            left,
            Inches(6.62),
            Inches(1.6),
            Inches(0.35),
            feat,
            font_size=9,
            color=WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )


def slide_challenges(prs):
    """Slide 12: Challenges & solutions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(
        slide, "Challenges Solved", "Real engineering problems encountered and resolved"
    )

    challenges = [
        (
            "Unrealistic Staffing Predictions",
            "Model extrapolated: predicted 4-5 agents when historical data showed 14+",
            "Added historical staffing floors (10th percentile of quality-meeting shifts) "
            "as lower bounds. Prevents model from predicting below observed reality.",
            ACCENT_RED,
            SUCCESS,
        ),
        (
            "English Agent Scarcity",
            "8 English agents total, optimizer demanded 8/day — impossible with 6 leave days",
            "Added feasibility caps: available_english = max(1, int(8 × 0.8) - 1) = 5. "
            "Prevents infeasible staffing requirements that would block the scheduler.",
            ACCENT_RED,
            SUCCESS,
        ),
        (
            "Preference vs. Fairness Tension",
            "Adding preference soft constraint → solver trade-offs night shift equity",
            "Sparsity-adjusted Gini scoring (only ~80 night slots for 50 agents), "
            "tunable weight balance. Accept mathematical limits of sparse distributions.",
            AMBER,
            SUCCESS,
        ),
        (
            "Windows Encoding",
            "Emoji in CLI output crashed on Windows cp1252 console encoding",
            "Wrapped stdout with io.TextIOWrapper(encoding='utf-8', errors='replace'). "
            "Cross-platform compatible.",
            AMBER,
            SUCCESS,
        ),
    ]

    for i, (title, problem, solution, prob_color, sol_color) in enumerate(challenges):
        row = i // 2
        col = i % 2
        left = Inches(0.6 + col * 6.3)
        top = Inches(1.6 + row * 2.7)

        add_shape(slide, left, top, Inches(5.9), Inches(2.3), LIGHT_GRAY, 0.03)
        add_textbox(
            slide,
            left + Inches(0.2),
            top + Inches(0.1),
            Inches(5.5),
            Inches(0.35),
            title,
            font_size=15,
            color=DARK,
            bold=True,
        )
        add_textbox(
            slide,
            left + Inches(0.2),
            top + Inches(0.5),
            Inches(5.5),
            Inches(0.55),
            f"Problem: {problem}",
            font_size=11,
            color=prob_color,
        )
        add_textbox(
            slide,
            left + Inches(0.2),
            top + Inches(1.15),
            Inches(5.5),
            Inches(0.95),
            f"Solution: {solution}",
            font_size=11,
            color=DARK,
        )


def slide_next_steps(prs):
    """Slide 13: Future improvements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(
        slide, "What's Next", "Production roadmap if this goes to production"
    )

    # Three columns
    categories = [
        (
            "Short Term",
            CYAN,
            [
                "Add real agent preference survey integration",
                "Auto-tune solver time based on solution gap",
                "A/B test scheduling configs",
                "Add shift swap request workflow",
            ],
        ),
        (
            "Medium Term",
            DEEP_BLUE,
            [
                "Online learning: retrain models monthly on actual data",
                "Multi-month scheduling (rollover constraints)",
                "Agent skill-based routing (beyond Sr/Jr/Eng)",
                "Real-time re-scheduling on sick days",
            ],
        ),
        (
            "Long Term",
            DARK,
            [
                "Reinforcement learning for adaptive scheduling",
                "Integration with HR/payroll systems",
                "Multi-site scheduling optimization",
                "Automated anomaly detection in CSAT trends",
            ],
        ),
    ]

    for i, (title, color, items) in enumerate(categories):
        left = Inches(0.6 + i * 4.2)
        add_shape(slide, left, Inches(1.6), Inches(3.8), Inches(0.5), color, 0.05)
        add_textbox(
            slide,
            left,
            Inches(1.62),
            Inches(3.8),
            Inches(0.45),
            title,
            font_size=18,
            color=WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_bullet_list(
            slide,
            left + Inches(0.2),
            Inches(2.3),
            Inches(3.4),
            Inches(3.5),
            items,
            font_size=13,
        )


def slide_thank_you(prs):
    """Slide 14: Thank you / Q&A."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CYAN)

    add_textbox(
        slide,
        Inches(1),
        Inches(2.0),
        Inches(11),
        Inches(1.0),
        "Thank You",
        font_size=52,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        Inches(1),
        Inches(3.2),
        Inches(11),
        Inches(0.6),
        "Questions & Discussion",
        font_size=24,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    # Divider
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.0), Inches(2.3), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()

    # Repo info
    add_textbox(
        slide,
        Inches(1),
        Inches(4.5),
        Inches(11),
        Inches(0.4),
        "github.com/ougrid/arwn-aifse",
        font_size=16,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        Inches(1),
        Inches(5.2),
        Inches(11),
        Inches(0.4),
        "Python 3.11  ·  scikit-learn  ·  OR-Tools  ·  Streamlit  ·  55+ Tests  ·  13 Commits",
        font_size=13,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    # Brand
    add_textbox(
        slide,
        Inches(1),
        Inches(6.3),
        Inches(11),
        Inches(0.4),
        "Aura Bangkok Clinic",
        font_size=14,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )


# ── Main ──────────────────────────────────────────────────────────────


def generate_presentation(output_path="presentation_deck.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)  # 1. Title
    slide_problem(prs)  # 2. Problem statement
    slide_approach(prs)  # 3. Solution architecture
    slide_forecasting(prs)  # 4. Task A: Forecasting
    slide_scheduling(prs)  # 5. Task B: Scheduling
    slide_results(prs)  # 6. Results
    slide_decisions(prs)  # 7. Design decisions
    slide_bonus_preferences(prs)  # 8. Bonus: Preferences
    slide_bonus_fairness(prs)  # 9. Bonus: Fairness & Remediation
    slide_production(prs)  # 10. Production readiness
    slide_dashboard(prs)  # 11. Dashboard
    slide_challenges(prs)  # 12. Challenges
    slide_next_steps(prs)  # 13. What's next
    slide_thank_you(prs)  # 14. Thank you

    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"  {len(prs.slides)} slides, 16:9 widescreen")
    return output_path


if __name__ == "__main__":
    generate_presentation()
